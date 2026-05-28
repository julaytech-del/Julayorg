"""
Julay.org — User Catalog PowerPoint Generator
Generates a polished Arabic/English user guide for the platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ──────────────────────────────────────────────────────────────
PRIMARY    = RGBColor(0x63, 0x66, 0xF1)   # Indigo
SECONDARY  = RGBColor(0x8B, 0x5C, 0xF6)   # Violet
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # Dark navy
BG_LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)   # Near white
TEXT_DARK  = RGBColor(0x0F, 0x17, 0x2A)
TEXT_MED   = RGBColor(0x47, 0x55, 0x69)
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)
RED        = RGBColor(0xEF, 0x44, 0x44)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
SURFACE    = RGBColor(0xF1, 0xF5, 0xF9)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0), radius=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=Pt(12), bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_badge(slide, text, x, y, w, h, bg, fg=TEXT_LIGHT, font_size=Pt(9)):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, text, x, y, w, h,
             font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)

# ── Slide builders ────────────────────────────────────────────────────────────

def slide_cover():
    """Slide 1 — Cover"""
    sl = prs.slides.add_slide(blank_layout)

    # Full dark background
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_DARK)

    # Decorative gradient orb (simulated with overlapping shapes)
    orb_colors = [
        (RGBColor(0x63,0x66,0xF1), Inches(4),   Inches(1),   Inches(3.5), Inches(3.5)),
        (RGBColor(0x8B,0x5C,0xF6), Inches(8),   Inches(3),   Inches(3),   Inches(3)),
    ]
    for c, ox, oy, ow, oh in orb_colors:
        orb = sl.shapes.add_shape(9, ox, oy, ow, oh)  # oval
        orb.fill.solid()
        orb.fill.fore_color.rgb = c
        orb.line.fill.background()
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        sp = orb._element
        sp.attrib['style'] = 'opacity:0.12'

    # Logo box
    logo_box = sl.shapes.add_shape(1, Inches(1.2), Inches(1.8), Inches(0.7), Inches(0.7))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = PRIMARY
    logo_box.line.fill.background()
    add_text(sl, "J", Inches(1.2), Inches(1.78), Inches(0.7), Inches(0.7),
             font_size=Pt(28), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Product name
    add_text(sl, "Julay.org", Inches(1.2), Inches(2.6), Inches(5), Inches(0.8),
             font_size=Pt(52), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.LEFT)

    # Tagline Arabic
    add_text(sl, "دليل المستخدم الشامل", Inches(1.2), Inches(3.4), Inches(7), Inches(0.6),
             font_size=Pt(26), bold=False, color=RGBColor(0xA5,0xB4,0xFC), align=PP_ALIGN.LEFT)

    # Sub tagline
    add_text(sl, "منصة إدارة المشاريع المدعومة بالذكاء الاصطناعي",
             Inches(1.2), Inches(4.05), Inches(8), Inches(0.5),
             font_size=Pt(16), color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.LEFT)

    # Bottom pill badges
    badges = [("🤖  ذكاء اصطناعي", Inches(1.2)), ("📋  إدارة مهام", Inches(3.2)), ("📊  تقارير", Inches(5.2)), ("⚡  أتمتة", Inches(6.9))]
    for label, bx in badges:
        add_rect(sl, bx, Inches(5.3), Inches(1.7), Inches(0.42), fill=RGBColor(0x1E,0x29,0x3B))
        add_text(sl, label, bx + Inches(0.05), Inches(5.3), Inches(1.6), Inches(0.42),
                 font_size=Pt(11), color=RGBColor(0xA5,0xB4,0xFC), align=PP_ALIGN.CENTER)

    # Right side decorative card
    card_x = Inches(8.8)
    add_rect(sl, card_x, Inches(1.5), Inches(3.8), Inches(4.5),
             fill=RGBColor(0x1E,0x29,0x3B), line_color=RGBColor(0x33,0x41,0x55))
    add_text(sl, "✦  AI Studio", card_x + Inches(0.2), Inches(1.8), Inches(3.4), Inches(0.4),
             font_size=Pt(12), bold=True, color=RGBColor(0xA5,0xB4,0xFC))
    add_text(sl, 'اكتب وصف مشروعك...', card_x + Inches(0.2), Inches(2.25), Inches(3.4), Inches(0.35),
             font_size=Pt(10), color=RGBColor(0x64,0x74,0x8B), italic=True)
    tasks_demo = [
        ("✓  إعداد قاعدة البيانات",    GREEN),
        ("◎  تصميم واجهة المستخدم",  PRIMARY),
        ("◌  كتابة الوثائق",           GRAY),
    ]
    for i, (t, c) in enumerate(tasks_demo):
        add_rect(sl, card_x + Inches(0.2), Inches(2.75) + Inches(i*0.55),
                 Inches(3.4), Inches(0.44), fill=RGBColor(0x0F,0x17,0x2A),
                 line_color=RGBColor(0x33,0x41,0x55))
        add_text(sl, t, card_x + Inches(0.35), Inches(2.78) + Inches(i*0.55),
                 Inches(3.1), Inches(0.38), font_size=Pt(10), color=c)

    # Version
    add_text(sl, "الإصدار 1.0  •  2026", Inches(1.2), Inches(6.6), Inches(5), Inches(0.4),
             font_size=Pt(10), color=RGBColor(0x47,0x55,0x69), align=PP_ALIGN.LEFT)

    return sl


def slide_toc():
    """Slide 2 — Table of Contents"""
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    # Header bar
    add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), fill=PRIMARY)
    add_text(sl, "📋  محتويات الدليل", Inches(0.4), Inches(0.25), Inches(10), Inches(0.6),
             font_size=Pt(26), bold=True, color=TEXT_LIGHT)
    add_text(sl, "كل ما تحتاج معرفته للعمل على Julay", Inches(0.4), Inches(0.82), Inches(10), Inches(0.35),
             font_size=Pt(12), color=RGBColor(0xC7,0xD2,0xFE))

    sections = [
        ("01", "الصفحة الرئيسية والتسجيل",   "كيف تبدأ وتُنشئ حسابك"),
        ("02", "لوحة التحكم",                  "نظرة عامة على إحصائياتك"),
        ("03", "المشاريع والمهام",              "إنشاء وإدارة المشاريع"),
        ("04", "لوحة Kanban",                  "تنظيم المهام بصرياً"),
        ("05", "Gantt & Sprint",               "التخطيط الزمني والسبرينت"),
        ("06", "الذكاء الاصطناعي AI Studio",  "أدوات AI لتسريع العمل"),
        ("07", "التقويم وتوزيع العمل",         "التنظيم الزمني للفريق"),
        ("08", "تتبع الوقت",                   "قياس ساعات العمل"),
        ("09", "التقارير والإحصائيات",         "تحليل أداء الفريق"),
        ("10", "الأتمتة والإشعارات",           "قواعد تعمل تلقائياً"),
        ("11", "إدارة الفريق والإعدادات",      "الأعضاء والصلاحيات"),
        ("12", "خطط الاشتراك",                 "المقارنة بين الخطط"),
    ]

    cols = 2
    col_w = Inches(6.2)
    for i, (num, title, desc) in enumerate(sections):
        col = i % cols
        row = i // cols
        x = Inches(0.35) + col * col_w
        y = Inches(1.35) + row * Inches(0.73)

        add_rect(sl, x, y, col_w - Inches(0.2), Inches(0.62),
                 fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))

        # Number badge
        add_rect(sl, x + Inches(0.08), y + Inches(0.08), Inches(0.46), Inches(0.46), fill=PRIMARY)
        add_text(sl, num, x + Inches(0.08), y + Inches(0.06), Inches(0.46), Inches(0.46),
                 font_size=Pt(11), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

        add_text(sl, title, x + Inches(0.65), y + Inches(0.04), col_w - Inches(0.9), Inches(0.3),
                 font_size=Pt(12), bold=True, color=TEXT_DARK)
        add_text(sl, desc, x + Inches(0.65), y + Inches(0.32), col_w - Inches(0.9), Inches(0.25),
                 font_size=Pt(9.5), color=TEXT_MED)

    return sl


def feature_slide(icon, title, subtitle, features, tip=None, accent=PRIMARY):
    """Generic feature slide with bullet points."""
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    # Left panel
    add_rect(sl, 0, 0, Inches(4.5), SLIDE_H, fill=BG_DARK)

    # Icon circle
    icon_shape = sl.shapes.add_shape(9, Inches(0.6), Inches(0.8), Inches(1.2), Inches(1.2))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = accent
    icon_shape.line.fill.background()
    add_text(sl, icon, Inches(0.6), Inches(0.78), Inches(1.2), Inches(1.2),
             font_size=Pt(32), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Title
    add_text(sl, title, Inches(0.35), Inches(2.2), Inches(3.8), Inches(0.8),
             font_size=Pt(26), bold=True, color=TEXT_LIGHT)

    # Subtitle
    add_text(sl, subtitle, Inches(0.35), Inches(3.05), Inches(3.8), Inches(1.5),
             font_size=Pt(13), color=RGBColor(0x94,0xA3,0xB8), wrap=True)

    # Tip box on left panel
    if tip:
        add_rect(sl, Inches(0.25), Inches(5.5), Inches(4.0), Inches(1.3),
                 fill=RGBColor(0x1E,0x29,0x3B), line_color=accent, line_width=Pt(1))
        add_text(sl, "💡  نصيحة", Inches(0.4), Inches(5.55), Inches(3.7), Inches(0.3),
                 font_size=Pt(9.5), bold=True, color=accent)
        add_text(sl, tip, Inches(0.4), Inches(5.85), Inches(3.7), Inches(0.85),
                 font_size=Pt(10), color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

    # Right side: feature cards
    right_x = Inches(4.8)
    for i, (feat_icon, feat_title, feat_desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = right_x + col * Inches(4.1)
        y = Inches(0.7) + row * Inches(1.55)
        w = Inches(3.9)
        h = Inches(1.4)

        add_rect(sl, x, y, w, h, fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))

        # Icon circle small
        ic = sl.shapes.add_shape(9, x + Inches(0.15), y + Inches(0.15), Inches(0.55), Inches(0.55))
        ic.fill.solid()
        ic.fill.fore_color.rgb = RGBColor(0xEE,0xF2,0xFF)
        ic.line.fill.background()
        add_text(sl, feat_icon, x + Inches(0.15), y + Inches(0.12), Inches(0.55), Inches(0.55),
                 font_size=Pt(18), align=PP_ALIGN.CENTER)

        add_text(sl, feat_title, x + Inches(0.82), y + Inches(0.12), w - Inches(1.0), Inches(0.35),
                 font_size=Pt(12), bold=True, color=TEXT_DARK)
        add_text(sl, feat_desc, x + Inches(0.15), y + Inches(0.65), w - Inches(0.3), Inches(0.65),
                 font_size=Pt(10), color=TEXT_MED, wrap=True)

    return sl


def slide_dashboard():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    # Header
    add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), fill=BG_DARK)
    add_text(sl, "📊", Inches(0.35), Inches(0.2), Inches(0.7), Inches(0.7), font_size=Pt(28))
    add_text(sl, "لوحة التحكم الرئيسية", Inches(1.05), Inches(0.25), Inches(8), Inches(0.55),
             font_size=Pt(24), bold=True, color=TEXT_LIGHT)
    add_text(sl, "أول شيء تشوفه بعد الدخول — نظرة شاملة على كل شيء دفعة واحدة",
             Inches(1.05), Inches(0.78), Inches(10), Inches(0.3),
             font_size=Pt(11), color=RGBColor(0x94,0xA3,0xB8))

    # Stat cards row
    stats = [
        ("المشاريع الكلية",   "12",  PRIMARY, "📁"),
        ("المهام النشطة",     "48",  ORANGE,  "✅"),
        ("أعضاء الفريق",     "8",   GREEN,   "👥"),
        ("تم إنجازها",        "134", SECONDARY,"⭐"),
    ]
    for i, (label, val, color, ic) in enumerate(stats):
        x = Inches(0.3) + i * Inches(3.25)
        y = Inches(1.25)
        add_rect(sl, x, y, Inches(3.05), Inches(1.3), fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))
        # color bar top
        add_rect(sl, x, y, Inches(3.05), Inches(0.08), fill=color)
        add_text(sl, ic, x + Inches(2.45), y + Inches(0.2), Inches(0.5), Inches(0.5), font_size=Pt(22))
        add_text(sl, val, x + Inches(0.2), y + Inches(0.15), Inches(1.8), Inches(0.6),
                 font_size=Pt(36), bold=True, color=TEXT_DARK)
        add_text(sl, label, x + Inches(0.2), y + Inches(0.78), Inches(2.4), Inches(0.3),
                 font_size=Pt(10), color=TEXT_MED)

    # Chart placeholder left
    add_rect(sl, Inches(0.3), Inches(2.75), Inches(6.2), Inches(3.8),
             fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))
    add_text(sl, "📈  توزيع المهام الأسبوعي", Inches(0.5), Inches(2.85), Inches(5.5), Inches(0.4),
             font_size=Pt(12), bold=True, color=TEXT_DARK)

    # Fake bar chart
    bar_data = [30, 55, 42, 70, 48, 85, 60]
    days = ["الأح", "الإث", "الثل", "الأر", "الخم", "الجم", "السب"]
    for i, (val, day) in enumerate(zip(bar_data, days)):
        bx = Inches(0.6) + i * Inches(0.8)
        bh = Inches(val / 100 * 2.2)
        by = Inches(2.75) + Inches(2.7) - bh + Inches(0.6)
        add_rect(sl, bx, by, Inches(0.5), bh, fill=PRIMARY)
        add_text(sl, day, bx, Inches(5.6), Inches(0.5), Inches(0.3),
                 font_size=Pt(8), color=TEXT_MED, align=PP_ALIGN.CENTER)

    # Recent activity right
    add_rect(sl, Inches(6.75), Inches(2.75), Inches(6.25), Inches(3.8),
             fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))
    add_text(sl, "⚡  آخر النشاطات", Inches(6.95), Inches(2.85), Inches(5.5), Inches(0.4),
             font_size=Pt(12), bold=True, color=TEXT_DARK)

    activities = [
        ("✦", PRIMARY,   "أحمد أنشأ مهمة جديدة",         "منذ 5 دقائق"),
        ("✓", GREEN,     "سارة أنجزت 'تصميم الشعار'",     "منذ 20 دقيقة"),
        ("★", SECONDARY, "AI ولّد خطة مشروع التطبيق",     "منذ ساعة"),
        ("↻", ORANGE,    "خالد غيّر حالة مهمة إلى Review","منذ ساعتين"),
        ("→", RGBColor(0x0E,0xA5,0xE9), "عُيّنت مهمة لمحمد", "أمس"),
    ]
    for i, (ic, col, text, time) in enumerate(activities):
        y = Inches(3.35) + i * Inches(0.55)
        add_rect(sl, Inches(6.95), y, Inches(0.3), Inches(0.3), fill=col)
        add_text(sl, ic, Inches(6.95), y - Inches(0.02), Inches(0.3), Inches(0.34),
                 font_size=Pt(12), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_text(sl, text, Inches(7.35), y + Inches(0.02), Inches(4.4), Inches(0.26),
                 font_size=Pt(10), color=TEXT_DARK)
        add_text(sl, time, Inches(11.2), y + Inches(0.04), Inches(1.6), Inches(0.22),
                 font_size=Pt(9), color=GRAY, align=PP_ALIGN.RIGHT)

    return sl


def slide_projects():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), fill=RGBColor(0x06,0x62,0x34) if False else PRIMARY)
    add_text(sl, "📁", Inches(0.35), Inches(0.2), Inches(0.7), Inches(0.7), font_size=Pt(28))
    add_text(sl, "المشاريع والمهام", Inches(1.05), Inches(0.25), Inches(8), Inches(0.55),
             font_size=Pt(24), bold=True, color=TEXT_LIGHT)
    add_text(sl, "أنشئ مشروعاً، أضف المهام، وعيّنها للفريق في دقائق",
             Inches(1.05), Inches(0.78), Inches(10), Inches(0.3),
             font_size=Pt(11), color=RGBColor(0x94,0xA3,0xB8))

    # Project cards
    projects = [
        ("تطبيق التوصيل",       "in_progress", "78%", "4 مهام متبقية",   "#4F46E5"),
        ("موقع الشركة",          "review",       "92%", "1 مهمة متبقية",   "#8B5CF6"),
        ("نظام الفواتير",         "planned",     "15%", "12 مهمة متبقية",  "#0EA5E9"),
    ]
    for i, (name, status, pct, sub, color) in enumerate(projects):
        x = Inches(0.3) + i * Inches(4.35)
        y = Inches(1.35)
        w = Inches(4.1)

        add_rect(sl, x, y, w, Inches(2.5), fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))
        # Color bar
        add_rect(sl, x, y, w, Inches(0.09), fill=RGBColor.from_string(color[1:]))

        add_text(sl, name, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.4),
                 font_size=Pt(14), bold=True, color=TEXT_DARK)

        status_colors = {"in_progress": (PRIMARY, "قيد التنفيذ"), "review": (ORANGE, "مراجعة"), "planned": (GRAY, "مخطط")}
        sc, sl_label = status_colors[status]
        add_rect(sl, x + Inches(0.2), y + Inches(0.68), Inches(1.3), Inches(0.32), fill=RGBColor(0xEE,0xF2,0xFF))
        add_text(sl, sl_label, x + Inches(0.2), y + Inches(0.68), Inches(1.3), Inches(0.32),
                 font_size=Pt(9.5), bold=True, color=sc, align=PP_ALIGN.CENTER)

        # Progress bar track
        add_rect(sl, x + Inches(0.2), y + Inches(1.2), w - Inches(0.4), Inches(0.12),
                 fill=RGBColor(0xE2,0xE8,0xF0))
        pct_val = int(pct.replace('%','')) / 100
        add_rect(sl, x + Inches(0.2), y + Inches(1.2),
                 (w - Inches(0.4)) * pct_val, Inches(0.12), fill=RGBColor.from_string(color[1:]))

        add_text(sl, pct, x + w - Inches(0.6), y + Inches(1.0), Inches(0.5), Inches(0.25),
                 font_size=Pt(10), bold=True, color=TEXT_MED, align=PP_ALIGN.RIGHT)
        add_text(sl, sub, x + Inches(0.2), y + Inches(1.42), w - Inches(0.4), Inches(0.25),
                 font_size=Pt(9.5), color=TEXT_MED)

        # Member avatars (circles)
        for j in range(3):
            av = sl.shapes.add_shape(9, x + Inches(0.2) + j*Inches(0.38), y + Inches(1.85),
                                      Inches(0.32), Inches(0.32))
            av.fill.solid()
            av.fill.fore_color.rgb = [PRIMARY, SECONDARY, GREEN][j]
            av.line.fill.background()

    # How-to steps
    add_text(sl, "كيف تُنشئ مشروعاً جديداً؟", Inches(0.3), Inches(4.05), Inches(12), Inches(0.4),
             font_size=Pt(14), bold=True, color=TEXT_DARK)

    steps = [
        ("1", "اضغط 'مشروع جديد'",     "من قائمة المشاريع"),
        ("2", "أدخل اسم المشروع",       "ووصفاً مختصراً"),
        ("3", "حدد الفريق والتواريخ",   "ابتداء وانتهاء"),
        ("4", "ابدأ إضافة المهام",       "يدوياً أو بالـ AI"),
    ]
    for i, (num, title, sub) in enumerate(steps):
        x = Inches(0.3) + i * Inches(3.25)
        y = Inches(4.55)
        # Arrow connector
        if i > 0:
            add_text(sl, "→", x - Inches(0.35), y + Inches(0.28), Inches(0.3), Inches(0.4),
                     font_size=Pt(16), color=PRIMARY, align=PP_ALIGN.CENTER)

        add_rect(sl, x, y, Inches(3.05), Inches(1.1), fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))
        # step number
        add_rect(sl, x + Inches(0.15), y + Inches(0.15), Inches(0.38), Inches(0.38), fill=PRIMARY)
        add_text(sl, num, x + Inches(0.15), y + Inches(0.13), Inches(0.38), Inches(0.38),
                 font_size=Pt(13), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_text(sl, title, x + Inches(0.65), y + Inches(0.12), Inches(2.2), Inches(0.3),
                 font_size=Pt(11), bold=True, color=TEXT_DARK)
        add_text(sl, sub, x + Inches(0.15), y + Inches(0.58), Inches(2.8), Inches(0.35),
                 font_size=Pt(10), color=TEXT_MED)

    # Tip
    add_rect(sl, Inches(0.3), Inches(5.85), Inches(12.7), Inches(0.85),
             fill=RGBColor(0xEE,0xF2,0xFF), line_color=PRIMARY, line_width=Pt(1))
    add_text(sl, "💡  يمكنك استخدام الذكاء الاصطناعي لتوليد مهام المشروع تلقائياً — فقط اكتب وصفاً للمشروع واضغط 'توليد AI'",
             Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.65),
             font_size=Pt(11), color=PRIMARY, wrap=True)

    return sl


def slide_kanban():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    add_rect(sl, 0, 0, SLIDE_W, Inches(1.0), fill=RGBColor(0x04,0x46,0x87) if False else SECONDARY)
    add_text(sl, "🗂️", Inches(0.35), Inches(0.18), Inches(0.7), Inches(0.65), font_size=Pt(26))
    add_text(sl, "لوحة Kanban", Inches(1.05), Inches(0.2), Inches(8), Inches(0.5),
             font_size=Pt(24), bold=True, color=TEXT_LIGHT)
    add_text(sl, "شاهد كل مهام مشروعك مرتبة في أعمدة حسب حالتها — اسحب وأفلت لتغيير الحالة",
             Inches(1.05), Inches(0.72), Inches(10), Inches(0.28),
             font_size=Pt(11), color=RGBColor(0xDD, 0xD6, 0xFE))

    columns = [
        ("📥 Backlog",    GRAY,      [("بحث المنافسين", "منخفض"), ("توثيق API", "متوسط")]),
        ("📋 To Do",      PRIMARY,   [("تصميم الشعار", "عالي"), ("إعداد السيرفر", "حرج")]),
        ("⚙️ قيد التنفيذ", ORANGE,  [("تطوير الواجهة", "عالي")]),
        ("👁️ مراجعة",    SECONDARY, [("اختبار الأداء", "متوسط"), ("Code Review", "عالي")]),
        ("✅ منجز",       GREEN,     [("قاعدة البيانات", "متوسط"), ("التصاميم", "منخفض")]),
    ]

    prio_colors = {"منخفض": GRAY, "متوسط": ORANGE, "عالي": RED, "حرج": RGBColor(0xDC,0x26,0x26)}
    col_w = Inches(2.45)

    for ci, (col_name, col_color, tasks) in enumerate(columns):
        cx = Inches(0.25) + ci * (col_w + Inches(0.12))

        # Column header
        add_rect(sl, cx, Inches(1.1), col_w, Inches(0.45), fill=col_color)
        add_text(sl, col_name, cx + Inches(0.1), Inches(1.12), col_w - Inches(0.2), Inches(0.4),
                 font_size=Pt(11), bold=True, color=TEXT_LIGHT)

        # Column body
        add_rect(sl, cx, Inches(1.55), col_w, Inches(5.65),
                 fill=RGBColor(0xF1,0xF5,0xF9), line_color=BORDER, line_width=Pt(0.5))

        # Task cards
        for ti, (task_name, prio) in enumerate(tasks):
            ty = Inches(1.65) + ti * Inches(1.45)
            add_rect(sl, cx + Inches(0.1), ty, col_w - Inches(0.2), Inches(1.3),
                     fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))

            add_text(sl, task_name, cx + Inches(0.2), ty + Inches(0.12), col_w - Inches(0.4), Inches(0.35),
                     font_size=Pt(11), bold=True, color=TEXT_DARK, wrap=True)

            pc = prio_colors.get(prio, GRAY)
            add_rect(sl, cx + Inches(0.2), ty + Inches(0.55), Inches(0.1), Inches(0.1), fill=pc)
            add_text(sl, prio, cx + Inches(0.35), ty + Inches(0.5), Inches(1.0), Inches(0.22),
                     font_size=Pt(9), color=pc)

            # Member avatar
            av = sl.shapes.add_shape(9, cx + col_w - Inches(0.5), ty + Inches(0.92),
                                      Inches(0.3), Inches(0.3))
            av.fill.solid()
            av.fill.fore_color.rgb = PRIMARY
            av.line.fill.background()

        # + Add task button
        add_rect(sl, cx + Inches(0.1), Inches(6.9), col_w - Inches(0.2), Inches(0.3),
                 fill=RGBColor(0xE2,0xE8,0xF0))
        add_text(sl, "+ إضافة مهمة", cx + Inches(0.1), Inches(6.9), col_w - Inches(0.2), Inches(0.3),
                 font_size=Pt(9), color=TEXT_MED, align=PP_ALIGN.CENTER)

    return sl


def slide_ai():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_DARK)

    # Stars decoration
    for pos in [(Inches(0.5), Inches(0.3)), (Inches(12.5), Inches(1)), (Inches(6), Inches(7.1))]:
        add_text(sl, "✦", pos[0], pos[1], Inches(0.3), Inches(0.3),
                 font_size=Pt(10), color=RGBColor(0x33,0x41,0x55))

    # Header
    add_text(sl, "🤖  AI Studio", Inches(0.5), Inches(0.25), Inches(10), Inches(0.7),
             font_size=Pt(32), bold=True, color=TEXT_LIGHT)
    add_text(sl, "الذكاء الاصطناعي يعمل معك — وليس بدلاً منك",
             Inches(0.5), Inches(0.95), Inches(10), Inches(0.35),
             font_size=Pt(14), color=RGBColor(0xA5,0xB4,0xFC))

    # 4 AI tool cards
    tools = [
        ("🚀", "توليد خطة المشروع",    PRIMARY,
         "اكتب وصف مشروعك\n← AI ينشئ المهام والأولويات والجداول الزمنية تلقائياً",
         "مثال: 'تطبيق توصيل طعام لـ iOS'"),
        ("📋", "تقرير Standup اليومي", GREEN,
         "اختر المشروع\n← AI يُلخص الحالة ويُحدد نقاط الخطر",
         "يعطيك: 🟢 سليم / 🟡 تحذير / 🔴 خطر"),
        ("📈", "تحليل الأداء",          ORANGE,
         "اختر المشروع\n← AI يُعطيك نسبة أداء ويوصي بتحسينات",
         "نتيجة من 0 إلى 100% مع اقتراحات محددة"),
        ("🔄", "إعادة التخطيط",         SECONDARY,
         "أدخل سبب التغيير\n← AI يُعيد ترتيب المهام والجداول",
         "مثال: 'تأخر العميل في الموافقة'"),
    ]

    for i, (icon, title, color, desc, example) in enumerate(tools):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.55) + row * Inches(2.7)
        w = Inches(6.1)
        h = Inches(2.5)

        add_rect(sl, x, y, w, h, fill=RGBColor(0x1E,0x29,0x3B),
                 line_color=color, line_width=Pt(1.5))

        # Top color accent
        add_rect(sl, x, y, w, Inches(0.07), fill=color)

        add_text(sl, icon, x + Inches(0.2), y + Inches(0.18), Inches(0.6), Inches(0.6),
                 font_size=Pt(26))
        add_text(sl, title, x + Inches(0.85), y + Inches(0.2), w - Inches(1.1), Inches(0.4),
                 font_size=Pt(14), bold=True, color=TEXT_LIGHT)

        add_text(sl, desc, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.9),
                 font_size=Pt(11), color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

        add_rect(sl, x + Inches(0.2), y + Inches(1.82), w - Inches(0.4), Inches(0.5),
                 fill=RGBColor(0x0F,0x17,0x2A))
        add_text(sl, "💡 " + example, x + Inches(0.3), y + Inches(1.85), w - Inches(0.6), Inches(0.42),
                 font_size=Pt(9.5), color=color, italic=True, wrap=True)

    return sl


def slide_time_reports():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    # Split header
    add_rect(sl, 0, 0, Inches(6.665), Inches(1.0), fill=GREEN)
    add_rect(sl, Inches(6.665), 0, Inches(6.665), Inches(1.0), fill=RGBColor(0x0E,0x83,0x78) if False else ORANGE)

    add_text(sl, "⏱️  تتبع الوقت", Inches(0.3), Inches(0.22), Inches(5), Inches(0.55),
             font_size=Pt(22), bold=True, color=TEXT_LIGHT)
    add_text(sl, "📊  التقارير", Inches(7.0), Inches(0.22), Inches(5), Inches(0.55),
             font_size=Pt(22), bold=True, color=TEXT_LIGHT)

    # Divider
    add_rect(sl, Inches(6.615), Inches(1.0), Inches(0.1), Inches(6.5), fill=BORDER)

    # Time Tracking section
    # Timer display
    add_rect(sl, Inches(0.3), Inches(1.15), Inches(6.1), Inches(1.8),
             fill=BG_DARK, line_color=GREEN, line_width=Pt(1.5))
    add_text(sl, "01:24:37", Inches(0.5), Inches(1.25), Inches(4), Inches(0.9),
             font_size=Pt(42), bold=True, color=TEXT_LIGHT)
    add_text(sl, "تصميم واجهة المستخدم", Inches(0.5), Inches(2.1), Inches(4), Inches(0.35),
             font_size=Pt(11), color=GRAY)
    # Stop button
    add_rect(sl, Inches(4.8), Inches(1.45), Inches(1.4), Inches(0.55), fill=RED)
    add_text(sl, "⏹  إيقاف", Inches(4.8), Inches(1.47), Inches(1.4), Inches(0.5),
             font_size=Pt(12), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    add_text(sl, "سجلات اليوم:", Inches(0.3), Inches(3.15), Inches(5), Inches(0.3),
             font_size=Pt(11), bold=True, color=TEXT_DARK)

    time_entries = [
        ("تصميم الشعار",          "02:30",  True,  GREEN),
        ("اجتماع مع العميل",      "01:00",  True,  GREEN),
        ("مراجعة الكود",           "00:45",  False, GRAY),
    ]
    for i, (task, dur, billable, col) in enumerate(time_entries):
        y = Inches(3.55) + i * Inches(0.62)
        add_rect(sl, Inches(0.3), y, Inches(6.1), Inches(0.52),
                 fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.5))
        add_rect(sl, Inches(0.3), y, Inches(0.06), Inches(0.52), fill=col)
        add_text(sl, task, Inches(0.5), y + Inches(0.12), Inches(3.5), Inches(0.3),
                 font_size=Pt(11), color=TEXT_DARK)
        add_text(sl, dur, Inches(4.2), y + Inches(0.12), Inches(0.9), Inches(0.3),
                 font_size=Pt(13), bold=True, color=TEXT_MED, align=PP_ALIGN.CENTER)
        if billable:
            add_text(sl, "💰", Inches(5.3), y + Inches(0.12), Inches(0.4), Inches(0.3), font_size=Pt(14))

    add_text(sl, "إجمالي اليوم: 4h 15m", Inches(0.3), Inches(5.7), Inches(4), Inches(0.3),
             font_size=Pt(12), bold=True, color=GREEN)

    # Reports section
    add_text(sl, "أنواع التقارير المتاحة:", Inches(6.85), Inches(1.15), Inches(5.9), Inches(0.35),
             font_size=Pt(12), bold=True, color=TEXT_DARK)

    report_types = [
        ("📋", "تقرير المهام",           "حالة، أولوية، المُكلَّف"),
        ("📁", "تقرير المشاريع",          "معدلات الإنجاز والصحة"),
        ("👥", "أداء الفريق",             "إنتاجية كل عضو"),
        ("📅", "تقرير الجدول الزمني",    "المخطط مقابل الفعلي"),
    ]
    for i, (ic, title, sub) in enumerate(report_types):
        y = Inches(1.65) + i * Inches(0.78)
        add_rect(sl, Inches(6.85), y, Inches(6.1), Inches(0.65),
                 fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))
        add_text(sl, ic, Inches(6.95), y + Inches(0.1), Inches(0.5), Inches(0.45), font_size=Pt(22))
        add_text(sl, title, Inches(7.55), y + Inches(0.06), Inches(4), Inches(0.3),
                 font_size=Pt(12), bold=True, color=TEXT_DARK)
        add_text(sl, sub, Inches(7.55), y + Inches(0.35), Inches(4), Inches(0.25),
                 font_size=Pt(9.5), color=TEXT_MED)

    # Export options
    add_rect(sl, Inches(6.85), Inches(4.95), Inches(6.1), Inches(0.8),
             fill=RGBColor(0xEE,0xF2,0xFF), line_color=PRIMARY, line_width=Pt(1))
    add_text(sl, "تصدير التقارير:", Inches(7.0), Inches(5.0), Inches(2), Inches(0.3),
             font_size=Pt(11), bold=True, color=PRIMARY)
    for i, (label, x_off) in enumerate([("🖨️ طباعة", 0), ("📄 PDF", 1.6), ("📅 تحديد تاريخ", 3.2)]):
        add_rect(sl, Inches(7.0) + Inches(x_off), Inches(5.35), Inches(1.4), Inches(0.3),
                 fill=TEXT_LIGHT, line_color=PRIMARY, line_width=Pt(0.75))
        add_text(sl, label, Inches(7.0) + Inches(x_off), Inches(5.36), Inches(1.4), Inches(0.28),
                 font_size=Pt(9.5), color=PRIMARY, align=PP_ALIGN.CENTER)

    # Tip
    add_rect(sl, Inches(6.85), Inches(5.95), Inches(6.1), Inches(0.75),
             fill=RGBColor(0xEC,0xFD,0xF5), line_color=GREEN, line_width=Pt(1))
    add_text(sl, "💡  حدد المهام كـ 'Billable' لتتبع الساعات القابلة للفوترة وإرسالها للعميل",
             Inches(7.0), Inches(6.05), Inches(5.7), Inches(0.58),
             font_size=Pt(10), color=RGBColor(0x06,0x5F,0x46), wrap=True)

    return sl


def slide_automations():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    add_rect(sl, 0, 0, SLIDE_W, Inches(1.0), fill=ORANGE)
    add_text(sl, "⚡", Inches(0.35), Inches(0.18), Inches(0.65), Inches(0.65), font_size=Pt(28))
    add_text(sl, "الأتمتة — اعمل أقل، أنجز أكثر", Inches(1.05), Inches(0.2), Inches(9), Inches(0.5),
             font_size=Pt(24), bold=True, color=TEXT_LIGHT)
    add_text(sl, "حدد قواعد تعمل تلقائياً عند حدوث أي حدث في المشروع",
             Inches(1.05), Inches(0.72), Inches(9), Inches(0.28),
             font_size=Pt(11), color=RGBColor(0xFF,0xF7,0xED))

    # Formula visual
    add_rect(sl, Inches(0.3), Inches(1.15), Inches(12.7), Inches(0.85),
             fill=RGBColor(0xFF,0xF7,0xED), line_color=ORANGE, line_width=Pt(1))

    for label, x, color in [
        ("عند حدوث...\n(المحفز)", Inches(0.5), ORANGE),
        ("  →  ", Inches(3.8), TEXT_MED),
        ("نفّذ...\n(الإجراء)", Inches(5.0), PRIMARY),
        ("  +  شروط اختيارية", Inches(8.5), GRAY),
    ]:
        add_text(sl, label, x, Inches(1.22), Inches(2.5), Inches(0.65),
                 font_size=Pt(12), bold=(color != TEXT_MED and color != GRAY),
                 color=color, align=PP_ALIGN.CENTER if '→' in label else PP_ALIGN.LEFT)

    # Triggers column
    add_text(sl, "⚡  المحفزات المتاحة", Inches(0.3), Inches(2.2), Inches(5.5), Inches(0.38),
             font_size=Pt(13), bold=True, color=ORANGE)

    triggers = [
        ("إنشاء مهمة جديدة",          "task.created"),
        ("تغيير حالة المهمة",          "task.status_changed"),
        ("اقتراب موعد التسليم (48 ساعة)", "task.due_soon"),
        ("تأخر المهمة عن موعدها",       "task.overdue"),
        ("تعيين مهمة لشخص",            "task.assigned"),
        ("إنشاء مشروع جديد",           "project.created"),
    ]
    for i, (label, code) in enumerate(triggers):
        y = Inches(2.7) + i * Inches(0.65)
        add_rect(sl, Inches(0.3), y, Inches(5.8), Inches(0.52),
                 fill=RGBColor(0xFF,0xF7,0xED), line_color=ORANGE, line_width=Pt(0.5))
        add_text(sl, "⚡  " + label, Inches(0.45), y + Inches(0.1), Inches(3.8), Inches(0.32),
                 font_size=Pt(11), color=TEXT_DARK)
        add_text(sl, code, Inches(4.3), y + Inches(0.14), Inches(1.7), Inches(0.24),
                 font_size=Pt(8.5), color=GRAY, italic=True)

    # Actions column
    add_text(sl, "🎯  الإجراءات المتاحة", Inches(6.55), Inches(2.2), Inches(5.5), Inches(0.38),
             font_size=Pt(13), bold=True, color=PRIMARY)

    actions = [
        ("🔔", "إرسال إشعار لعضو",      "اختر العضو والرسالة"),
        ("🔄", "تغيير حالة المهمة",      "إلى: todo / in_progress / done..."),
        ("💬", "إضافة تعليق تلقائي",     "رسالة تُضاف للمهمة"),
        ("➕", "إنشاء مهمة فرعية",       "مهمة جديدة داخل المهمة"),
    ]
    for i, (ic, title, sub) in enumerate(actions):
        y = Inches(2.7) + i * Inches(0.9)
        add_rect(sl, Inches(6.55), y, Inches(6.45), Inches(0.75),
                 fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))
        add_text(sl, ic, Inches(6.65), y + Inches(0.12), Inches(0.45), Inches(0.5), font_size=Pt(22))
        add_text(sl, title, Inches(7.2), y + Inches(0.08), Inches(5), Inches(0.3),
                 font_size=Pt(12), bold=True, color=PRIMARY)
        add_text(sl, sub, Inches(7.2), y + Inches(0.4), Inches(5.5), Inches(0.25),
                 font_size=Pt(10), color=TEXT_MED)

    # Example rule
    add_rect(sl, Inches(0.3), Inches(6.25), Inches(12.7), Inches(0.9),
             fill=BG_DARK)
    add_text(sl, "📌  مثال عملي:", Inches(0.5), Inches(6.3), Inches(2.2), Inches(0.3),
             font_size=Pt(11), bold=True, color=ORANGE)
    add_text(sl, "عندما تتأخر مهمة ← أرسل إشعاراً لمدير المشروع + أضف تعليق 'المهمة تحتاج مراجعة'",
             Inches(2.8), Inches(6.32), Inches(10), Inches(0.55),
             font_size=Pt(11), color=TEXT_LIGHT, wrap=True)

    return sl


def slide_team_settings():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    # Split header
    add_rect(sl, 0, 0, Inches(6.665), Inches(1.0), fill=RGBColor(0x0E,0xA5,0xE9))
    add_rect(sl, Inches(6.665), 0, Inches(6.665), Inches(1.0), fill=RGBColor(0x47,0x55,0x69))

    add_text(sl, "👥  إدارة الفريق", Inches(0.3), Inches(0.22), Inches(5), Inches(0.55),
             font_size=Pt(22), bold=True, color=TEXT_LIGHT)
    add_text(sl, "⚙️  الإعدادات", Inches(7.0), Inches(0.22), Inches(5), Inches(0.55),
             font_size=Pt(22), bold=True, color=TEXT_LIGHT)

    add_rect(sl, Inches(6.615), Inches(1.0), Inches(0.1), Inches(6.5), fill=BORDER)

    # Team section
    add_text(sl, "أعضاء فريقك الحالي:", Inches(0.3), Inches(1.15), Inches(5.5), Inches(0.35),
             font_size=Pt(12), bold=True, color=TEXT_DARK)

    members = [
        ("أحمد محمد",     "مدير المشروع",  "admin",  PRIMARY),
        ("سارة علي",      "مصمم واجهات",   "member", GREEN),
        ("خالد السالم",   "مطور backend",  "member", SECONDARY),
        ("نورة الرشيد",   "QA Engineer",   "viewer", ORANGE),
    ]
    role_colors = {"admin": PRIMARY, "member": GREEN, "viewer": GRAY}
    role_labels = {"admin": "مدير", "member": "عضو", "viewer": "مشاهد"}

    for i, (name, role_title, role, color) in enumerate(members):
        y = Inches(1.6) + i * Inches(0.78)
        add_rect(sl, Inches(0.3), y, Inches(6.1), Inches(0.65),
                 fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.5))

        av = sl.shapes.add_shape(9, Inches(0.4), y + Inches(0.1), Inches(0.45), Inches(0.45))
        av.fill.solid(); av.fill.fore_color.rgb = color; av.line.fill.background()
        add_text(sl, name[0], Inches(0.4), y + Inches(0.08), Inches(0.45), Inches(0.45),
                 font_size=Pt(14), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

        add_text(sl, name, Inches(0.95), y + Inches(0.05), Inches(2.5), Inches(0.28),
                 font_size=Pt(11), bold=True, color=TEXT_DARK)
        add_text(sl, role_title, Inches(0.95), y + Inches(0.33), Inches(2.5), Inches(0.24),
                 font_size=Pt(9.5), color=TEXT_MED)

        rc = role_colors[role]
        add_rect(sl, Inches(4.0), y + Inches(0.17), Inches(1.1), Inches(0.28), fill=RGBColor(0xEE,0xF2,0xFF))
        add_text(sl, role_labels[role], Inches(4.0), y + Inches(0.17), Inches(1.1), Inches(0.28),
                 font_size=Pt(9.5), bold=True, color=rc, align=PP_ALIGN.CENTER)

    # Invite box
    add_rect(sl, Inches(0.3), Inches(4.85), Inches(6.1), Inches(0.85),
             fill=RGBColor(0xEE,0xF2,0xFF), line_color=PRIMARY, line_width=Pt(1))
    add_text(sl, "✉️  دعوة عضو جديد بالبريد الإلكتروني أو رابط الدعوة",
             Inches(0.5), Inches(4.95), Inches(5.7), Inches(0.55),
             font_size=Pt(11), color=PRIMARY, wrap=True)

    add_rect(sl, Inches(0.3), Inches(5.85), Inches(6.1), Inches(0.55),
             fill=BG_DARK)
    add_text(sl, "🔗  انسخ رابط الدعوة وأرسله لأي شخص — سيصله طلب الانضمام مباشرة",
             Inches(0.5), Inches(5.9), Inches(5.7), Inches(0.45),
             font_size=Pt(10), color=RGBColor(0xA5,0xB4,0xFC), wrap=True)

    # Settings section
    settings_tabs = [
        ("👤", "الملف الشخصي",    "الاسم، المسمى الوظيفي، الصورة، المنطقة الزمنية"),
        ("🔒", "الأمان",           "كلمة المرور، المصادقة الثنائية (2FA)"),
        ("🔔", "الإشعارات",        "اختر الإشعارات التي تريد استلامها"),
        ("💳", "الاشتراك",         "خطتك الحالية وخيارات الترقية"),
        ("🔗", "التكاملات",        "GitHub، Google Calendar، وأكثر"),
        ("🔔", "Webhooks",         "ربط المنصة بأدوات خارجية"),
    ]
    for i, (ic, tab, desc) in enumerate(settings_tabs):
        col = i % 2
        row = i // 2
        x = Inches(6.85) + col * Inches(3.2)
        y = Inches(1.15) + row * Inches(0.88)

        add_rect(sl, x, y, Inches(3.0), Inches(0.75),
                 fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.5))
        add_text(sl, ic, x + Inches(0.1), y + Inches(0.1), Inches(0.45), Inches(0.55), font_size=Pt(22))
        add_text(sl, tab, x + Inches(0.65), y + Inches(0.05), Inches(2.2), Inches(0.3),
                 font_size=Pt(11), bold=True, color=TEXT_DARK)
        add_text(sl, desc, x + Inches(0.65), y + Inches(0.35), Inches(2.2), Inches(0.3),
                 font_size=Pt(9), color=TEXT_MED, wrap=True)

    # Language note
    add_rect(sl, Inches(6.85), Inches(3.8), Inches(6.1), Inches(0.7),
             fill=RGBColor(0xEC,0xFD,0xF5), line_color=GREEN, line_width=Pt(1))
    add_text(sl, "🌍  اللغات: العربية • English • Français  —  تبديل فوري من إعدادات الملف الشخصي",
             Inches(7.0), Inches(3.9), Inches(5.7), Inches(0.5),
             font_size=Pt(10.5), color=RGBColor(0x06,0x5F,0x46), wrap=True)

    # 2FA note
    add_rect(sl, Inches(6.85), Inches(4.65), Inches(6.1), Inches(0.7),
             fill=RGBColor(0xFF,0xF7,0xED), line_color=ORANGE, line_width=Pt(1))
    add_text(sl, "🛡️  فعّل المصادقة الثنائية (2FA) لحماية إضافية — من الإعدادات ← الأمان",
             Inches(7.0), Inches(4.75), Inches(5.7), Inches(0.5),
             font_size=Pt(10.5), color=RGBColor(0x92,0x40,0x0E), wrap=True)

    return sl


def slide_pricing():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_DARK)

    add_text(sl, "💎  خطط الاشتراك", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
             font_size=Pt(30), bold=True, color=TEXT_LIGHT)
    add_text(sl, "اختر الخطة التي تناسب فريقك",
             Inches(0.5), Inches(0.82), Inches(10), Inches(0.35),
             font_size=Pt(13), color=RGBColor(0x94,0xA3,0xB8))

    plans = [
        {
            "name": "Free",
            "price": "مجاني",
            "desc": "للأفراد والفرق الصغيرة",
            "color": GRAY,
            "features": ["3 مشاريع", "3 أعضاء", "1 GB تخزين", "Kanban و List", "بدون AI"],
            "popular": False,
        },
        {
            "name": "Starter",
            "price": "$15 / شهر",
            "desc": "للفرق النامية",
            "color": RGBColor(0x0E,0xA5,0xE9),
            "features": ["10 مشاريع", "10 أعضاء", "10 GB", "30 طلب AI/شهر", "Gantt", "Time Tracking"],
            "popular": False,
        },
        {
            "name": "Professional",
            "price": "$47 / شهر",
            "desc": "للفرق الجادة",
            "color": PRIMARY,
            "features": ["مشاريع غير محدودة", "25 عضو", "50 GB", "500 طلب AI", "Automations", "Sprint & Workload", "Reports & Export"],
            "popular": True,
        },
        {
            "name": "Business",
            "price": "$79 / شهر",
            "desc": "للمؤسسات الكبيرة",
            "color": SECONDARY,
            "features": ["كل شيء غير محدود", "مستخدمون غير محدودون", "200 GB", "2,000 طلب AI", "API & Webhooks", "SSO / SAML", "دعم مخصص", "SLA 99.9%"],
            "popular": False,
        },
    ]

    plan_w = Inches(3.15)
    for i, plan in enumerate(plans):
        x = Inches(0.25) + i * (plan_w + Inches(0.12))
        y = Inches(1.3)

        # Card bg
        card_color = RGBColor(0x1E,0x29,0x3B) if not plan["popular"] else plan["color"]
        border_color = plan["color"]
        add_rect(sl, x, y, plan_w, Inches(5.85),
                 fill=card_color, line_color=border_color, line_width=Pt(2 if plan["popular"] else 1))

        if plan["popular"]:
            add_rect(sl, x, y, plan_w, Inches(0.38), fill=plan["color"])
            add_text(sl, "⭐  الأكثر شيوعاً", x, y + Inches(0.05), plan_w, Inches(0.28),
                     font_size=Pt(10), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
            y += Inches(0.38)

        add_text(sl, plan["name"], x + Inches(0.15), y + Inches(0.18), plan_w - Inches(0.3), Inches(0.42),
                 font_size=Pt(20), bold=True, color=plan["color"] if not plan["popular"] else TEXT_LIGHT)
        add_text(sl, plan["desc"], x + Inches(0.15), y + Inches(0.58), plan_w - Inches(0.3), Inches(0.3),
                 font_size=Pt(10), color=RGBColor(0x94,0xA3,0xB8))
        add_text(sl, plan["price"], x + Inches(0.15), y + Inches(0.92), plan_w - Inches(0.3), Inches(0.55),
                 font_size=Pt(22), bold=True, color=TEXT_LIGHT if plan["popular"] else plan["color"])

        # Divider
        add_rect(sl, x + Inches(0.15), y + Inches(1.55), plan_w - Inches(0.3), Inches(0.04),
                 fill=RGBColor(0x33,0x41,0x55))

        for fi, feat in enumerate(plan["features"]):
            fy = y + Inches(1.72) + fi * Inches(0.44)
            add_text(sl, "✓  " + feat, x + Inches(0.15), fy, plan_w - Inches(0.3), Inches(0.38),
                     font_size=Pt(10), color=TEXT_LIGHT if plan["popular"] else RGBColor(0xCB,0xD5,0xE1))

    # Annual discount note
    add_rect(sl, Inches(0.25), Inches(7.0), Inches(12.7), Inches(0.35),
             fill=RGBColor(0x1E,0x29,0x3B))
    add_text(sl, "💰  الدفع السنوي يوفر لك شهرين مجاناً  •  تجربة مجانية 14 يوم لخطط Starter و Professional",
             Inches(0.5), Inches(7.04), Inches(12.3), Inches(0.27),
             font_size=Pt(10), color=RGBColor(0xA5,0xB4,0xFC), align=PP_ALIGN.CENTER)

    return sl


def slide_extra_features():
    """Slide for Sprint, Portfolio, Workload, Calendar, Forms"""
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)

    add_rect(sl, 0, 0, SLIDE_W, Inches(1.0), fill=RGBColor(0x06,0x62,0x34) if False else SECONDARY)
    add_text(sl, "🔍", Inches(0.35), Inches(0.18), Inches(0.65), Inches(0.65), font_size=Pt(26))
    add_text(sl, "ميزات إضافية متقدمة", Inches(1.05), Inches(0.2), Inches(9), Inches(0.5),
             font_size=Pt(24), bold=True, color=TEXT_LIGHT)
    add_text(sl, "Sprint • Portfolio • Workload • Calendar • Forms",
             Inches(1.05), Inches(0.72), Inches(10), Inches(0.28),
             font_size=Pt(11), color=RGBColor(0xDDD6FE))

    extras = [
        ("🏃", "Sprint Board",
         "نظام السبرينت الأجايل",
         "قسّم عملك إلى دورات زمنية (عادة أسبوعين). اسحب المهام من Backlog للسبرينت الحالي وتابع التقدم باستمرار.",
         SECONDARY),
        ("📊", "Portfolio View",
         "عرض محفظة المشاريع",
         "شاهد كل مشاريعك في مكان واحد مع مؤشرات الصحة:\n🟢 On Track   🟡 At Risk   🔴 Off Track",
         RGBColor(0x0E,0xA5,0xE9)),
        ("📅", "Calendar View",
         "التقويم الزمني",
         "اعرض المهام ذات التواريخ على تقويم شهري/أسبوعي. انقر على أي يوم لإضافة مهمة مباشرة.",
         GREEN),
        ("⚖️", "Workload View",
         "توزيع العمل",
         "جدول أسبوعي يعرض ساعات كل عضو.\n🟢 أقل من 60%  🟡 60-90%  🔴 أكثر من 90%",
         ORANGE),
        ("📝", "Form Views",
         "النماذج المشتركة",
         "أنشئ نموذجاً مخصصاً وشاركه برابط عام — يمكن لأي شخص إرسال طلبات دون حساب.",
         RED),
        ("🔍", "Global Search",
         "البحث الشامل",
         "ابحث عن أي مهمة، مشروع، أو عضو في المنصة فوراً من شريط البحث العلوي.",
         RGBColor(0x14,0xB8,0xA6)),
    ]

    for i, (ic, title, subtitle, desc, color) in enumerate(extras):
        col = i % 3
        row = i // 3
        x = Inches(0.25) + col * Inches(4.35)
        y = Inches(1.15) + row * Inches(2.85)
        w = Inches(4.1)
        h = Inches(2.65)

        add_rect(sl, x, y, w, h, fill=TEXT_LIGHT, line_color=BORDER, line_width=Pt(0.75))
        add_rect(sl, x, y, w, Inches(0.08), fill=color)

        add_text(sl, ic, x + Inches(0.15), y + Inches(0.18), Inches(0.55), Inches(0.55), font_size=Pt(26))
        add_text(sl, title, x + Inches(0.82), y + Inches(0.15), w - Inches(1.0), Inches(0.35),
                 font_size=Pt(13), bold=True, color=TEXT_DARK)
        add_text(sl, subtitle, x + Inches(0.82), y + Inches(0.48), w - Inches(1.0), Inches(0.25),
                 font_size=Pt(9.5), color=color, bold=True)
        add_text(sl, desc, x + Inches(0.15), y + Inches(0.88), w - Inches(0.3), Inches(1.5),
                 font_size=Pt(10), color=TEXT_MED, wrap=True)

    return sl


def slide_closing():
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BG_DARK)

    # Decorative orbs
    orb = sl.shapes.add_shape(9, Inches(9), Inches(-0.5), Inches(5), Inches(5))
    orb.fill.solid(); orb.fill.fore_color.rgb = RGBColor(0x63,0x66,0xF1)
    orb.line.fill.background()

    orb2 = sl.shapes.add_shape(9, Inches(-1), Inches(5), Inches(4), Inches(4))
    orb2.fill.solid(); orb2.fill.fore_color.rgb = SECONDARY
    orb2.line.fill.background()

    add_text(sl, "J", Inches(1.5), Inches(1.5), Inches(1.1), Inches(1.1),
             font_size=Pt(52), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(1.5), Inches(1.5), Inches(1.1), Inches(1.1), fill=PRIMARY)
    add_text(sl, "J", Inches(1.5), Inches(1.48), Inches(1.1), Inches(1.1),
             font_size=Pt(48), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    add_text(sl, "ابدأ رحلتك مع Julay", Inches(1.2), Inches(2.85), Inches(9), Inches(0.85),
             font_size=Pt(40), bold=True, color=TEXT_LIGHT)

    add_text(sl, "كل ما تحتاجه لإدارة مشاريعك، فريقك، ووقتك — في مكان واحد.",
             Inches(1.2), Inches(3.78), Inches(9), Inches(0.5),
             font_size=Pt(17), color=RGBColor(0x94,0xA3,0xB8))

    # CTA button
    add_rect(sl, Inches(1.2), Inches(4.55), Inches(3.2), Inches(0.72),
             fill=PRIMARY)
    add_text(sl, "🚀  ابدأ مجاناً الآن", Inches(1.2), Inches(4.57), Inches(3.2), Inches(0.68),
             font_size=Pt(16), bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(4.6), Inches(4.55), Inches(2.5), Inches(0.72),
             fill=RGBColor(0x1E,0x29,0x3B), line_color=PRIMARY, line_width=Pt(1.5))
    add_text(sl, "julay.org", Inches(4.6), Inches(4.57), Inches(2.5), Inches(0.68),
             font_size=Pt(16), color=PRIMARY, align=PP_ALIGN.CENTER)

    # Quick reference
    add_rect(sl, Inches(1.2), Inches(5.55), Inches(9.5), Inches(1.4),
             fill=RGBColor(0x1E,0x29,0x3B), line_color=RGBColor(0x33,0x41,0x55))
    add_text(sl, "🔑  خلاصة سريعة:", Inches(1.4), Inches(5.65), Inches(3), Inches(0.3),
             font_size=Pt(11), bold=True, color=RGBColor(0xA5,0xB4,0xFC))

    tips = [
        "• لوحة التحكم: نظرة عامة فورية",
        "• Kanban: سحب وإفلات للحالات",
        "• AI Studio: وصف → مشروع كامل",
        "• الأتمتة: قواعد تعمل وحدها",
        "• التقارير: تصدير PDF بنقرة",
        "• الإعدادات: دعوة الفريق وإدارة الأدوار",
    ]
    for i, tip in enumerate(tips):
        col = i % 3
        row = i // 3
        x = Inches(1.4) + col * Inches(3.15)
        y = Inches(6.0) + row * Inches(0.42)
        add_text(sl, tip, x, y, Inches(3.0), Inches(0.38),
                 font_size=Pt(9.5), color=RGBColor(0xCB,0xD5,0xE1))

    return sl


# ── Build all slides ───────────────────────────────────────────────────────────

print("Building slides...")

slide_cover()
slide_toc()
slide_dashboard()
slide_projects()
slide_kanban()

feature_slide(
    "📅 / 🏃", "Gantt & Sprint Board",
    "خطط مشاريعك على محور الزمن\nوأدر عمل فريقك بدورات أجايل",
    [
        ("📅", "مخطط Gantt",        "شاهد المهام على محور الزمن، اسحب لتغيير التواريخ، وشاهد تبعيات المهام"),
        ("🏃", "Sprint Board",       "قسّم العمل إلى سبرينت 2 أسبوع، اسحب المهام من Backlog وتابع السرعة"),
        ("📊", "Story Points",       "قدّر جهد كل مهمة بالنقاط (= ساعات التقدير) لقياس سرعة الفريق"),
        ("⏱️", "أيام متبقية",       "عداد تنازلي لكل سبرينت مع مؤشر تقدم رسومي"),
    ],
    tip="Sprint مثالي للفرق التقنية — Gantt أفضل لمشاريع ذات مراحل محددة",
    accent=SECONDARY,
)

slide_ai()
slide_time_reports()
slide_automations()
slide_extra_features()
slide_team_settings()
slide_pricing()
slide_closing()

# ── Save ──────────────────────────────────────────────────────────────────────
output_path = "/Users/assi/elite-ai-factory/julayorg/Julay_User_Catalog.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
